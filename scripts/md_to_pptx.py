"""Convert slides.md to a basic PPTX.

Falls back from the rich reveal.js HTML to a minimal PPTX export, so users
who need editable PowerPoint can still get one. Each slide becomes one PPTX
slide; the first heading is the title, the rest is a single content
placeholder. Screenshots referenced by `![](path)` are embedded.

Usage:
    python md_to_pptx.py docs/training/action-tracker/slides.md out.pptx

Requires: pip install python-pptx
"""
from __future__ import annotations

import argparse
import re
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError:
    sys.exit("Install python-pptx: pip install python-pptx")

SLIDE_SEP = re.compile(r"^---\s*$", re.M)
HEADING = re.compile(r"^(#{1,3})\s+(.+)$", re.M)
IMAGE = re.compile(r"!\[[^\]]*\]\(([^)]+)\)")
COMMENT = re.compile(r"<!--.*?-->", re.S)
NOTE_SEP = re.compile(r"^Note:", re.M)


def parse_slide(raw: str) -> dict:
    raw = COMMENT.sub("", raw).strip()
    note = ""
    if NOTE_SEP.search(raw):
        body, note = NOTE_SEP.split(raw, 1)
    else:
        body = raw

    title = ""
    h = HEADING.search(body)
    if h:
        title = h.group(2).strip()
        body = HEADING.sub("", body, count=1)

    images = IMAGE.findall(body)
    text_body = IMAGE.sub("", body).strip()
    return {"title": title, "body": text_body, "images": images, "note": note.strip()}


def main() -> int:
    ap = argparse.ArgumentParser()
    ap.add_argument("slides_md")
    ap.add_argument("out_pptx")
    args = ap.parse_args()

    src = Path(args.slides_md)
    text = src.read_text(encoding="utf-8")
    slides = SLIDE_SEP.split(text)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    for raw in slides:
        s = parse_slide(raw)
        if not s["title"] and not s["body"] and not s["images"]:
            continue
        slide = prs.slides.add_slide(blank)

        # Title textbox
        if s["title"]:
            tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.3), Inches(0.9))
            tf = tb.text_frame
            tf.text = s["title"]
            tf.paragraphs[0].runs[0].font.size = Pt(36)
            tf.paragraphs[0].runs[0].font.bold = True

        # Image (first image only)
        if s["images"]:
            img_path = (src.parent / s["images"][0]).resolve()
            if img_path.exists():
                slide.shapes.add_picture(
                    str(img_path), Inches(0.5), Inches(1.5), height=Inches(5.5)
                )

        # Body text
        if s["body"]:
            tb = slide.shapes.add_textbox(
                Inches(7.0) if s["images"] else Inches(0.5),
                Inches(1.5) if s["images"] else Inches(1.5),
                Inches(5.8) if s["images"] else Inches(12.3),
                Inches(5.5),
            )
            tf = tb.text_frame
            tf.word_wrap = True
            for line in s["body"].splitlines():
                line = line.strip().rstrip("\\")
                if not line:
                    continue
                p = tf.add_paragraph()
                p.text = line.lstrip("- ").lstrip("* ").lstrip("0123456789. ")
                p.font.size = Pt(18)

        if s["note"]:
            slide.notes_slide.notes_text_frame.text = s["note"]

    prs.save(args.out_pptx)
    print(f"Wrote {args.out_pptx}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
